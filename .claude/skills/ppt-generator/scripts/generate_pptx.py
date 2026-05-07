#!/usr/bin/env python3
"""
Generate a .pptx file from a JSON slide definition.

Usage:
    python generate_pptx.py '<json_string>' '<output_path>'

JSON schema:
{
  "title": "Presentation Title",
  "author": "Author Name",           # optional
  "slides": [
    {
      "title": "Slide Title",
      "key_message": "...",
      "visual": "...",
      "notes": "..."                  # optional
    }
  ]
}
"""

import sys
import json

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("python-pptx not found. Installing...", file=sys.stderr)
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT     = RGBColor(0x16, 0x21, 0x3E)   # mid navy
HIGHLIGHT  = RGBColor(0x0F, 0x3E, 0x6F)   # bright navy accent bar
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW     = RGBColor(0xFF, 0xC3, 0x00)   # key-message accent


def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def _add_accent_bar(slide, top_inches=1.05):
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(top_inches), Inches(10), Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = HIGHLIGHT
    bar.line.fill.background()


def make_title_slide(prs, title, author=""):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, DARK_BG)

    # Big centred title
    _add_textbox(slide, title,
                 left=0.5, top=2.2, width=9.0, height=1.8,
                 font_size=40, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

    # Accent bar below title
    _add_accent_bar(slide, top_inches=3.9)

    # Author / subtitle
    if author:
        _add_textbox(slide, author,
                     left=0.5, top=4.1, width=9.0, height=0.5,
                     font_size=18, color=LIGHT_GRAY,
                     align=PP_ALIGN.CENTER)


def make_content_slide(prs, slide_data: dict):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, DARK_BG)

    title       = slide_data.get("title", "")
    key_message = slide_data.get("key_message", "")
    visual      = slide_data.get("visual", "")
    notes_text  = slide_data.get("notes", "")

    # Slide title
    _add_textbox(slide, title,
                 left=0.4, top=0.2, width=9.2, height=0.75,
                 font_size=28, bold=True, color=WHITE)

    _add_accent_bar(slide, top_inches=0.95)

    # Key message box (highlighted)
    km_box = slide.shapes.add_shape(
        1,
        Inches(0.4), Inches(1.1), Inches(9.2), Inches(0.9)
    )
    km_box.fill.solid()
    km_box.fill.fore_color.rgb = ACCENT
    km_box.line.fill.background()

    _add_textbox(slide,
                 f"Key message:  {key_message}",
                 left=0.55, top=1.12, width=9.0, height=0.86,
                 font_size=15, bold=False, color=YELLOW)

    # Visual instructions box
    vis_label_top = 2.15
    _add_textbox(slide, "VISUAL ASSET",
                 left=0.4, top=vis_label_top, width=2.0, height=0.3,
                 font_size=10, bold=True, color=HIGHLIGHT)

    vis_box = slide.shapes.add_shape(
        1,
        Inches(0.4), Inches(vis_label_top + 0.3),
        Inches(9.2), Inches(2.8)
    )
    vis_box.fill.solid()
    vis_box.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x22)
    vis_box.line.color.rgb = HIGHLIGHT

    _add_textbox(slide, visual,
                 left=0.55, top=vis_label_top + 0.35,
                 width=9.0, height=2.6,
                 font_size=14, color=LIGHT_GRAY)

    # Speaker notes
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text


def build_presentation(data: dict, output_path: str):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    make_title_slide(prs,
                     title=data.get("title", "Presentation"),
                     author=data.get("author", ""))

    for slide_data in data.get("slides", []):
        make_content_slide(prs, slide_data)

    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: generate_pptx.py '<json>' '<output_path>'", file=sys.stderr)
        sys.exit(1)

    raw_json    = sys.argv[1]
    output_path = sys.argv[2]

    try:
        data = json.loads(raw_json)
    except json.JSONDecodeError as e:
        print(f"Invalid JSON: {e}", file=sys.stderr)
        sys.exit(1)

    build_presentation(data, output_path)
